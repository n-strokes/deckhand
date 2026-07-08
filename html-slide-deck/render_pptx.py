# -*- coding: utf-8 -*-
"""Build a full-bleed PPTX rendition from a folder of slide screenshots.

The QUICK export: one PNG per slide, dropped full-bleed onto a blank 16:9 slide.
The result is a faithful picture of the HTML deck, not editable shapes. (For the
FULL export - every element rebuilt as native PPT objects - write a bespoke
python-pptx builder per the deck; see SKILL.md. Do not use this script for that.)

Usage:
    python3 render_pptx.py SHOTS_DIR --out NAME_html_rendition.pptx

SHOTS_DIR holds slideNN.png files; they are inserted in sorted filename order.

macOS note: rewritten from the original Windows PowerPoint-COM version to use
python-pptx, so it is cross-platform and needs no PowerPoint install.
Install the dependency once:  pip3 install python-pptx
"""
import argparse, glob, os, sys
from pptx import Presentation
from pptx.util import Inches

W_IN, H_IN = 13.333, 7.5            # widescreen 16:9
BLANK_LAYOUT = 6                    # the blank layout in the default template

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("shots", help="folder of slideNN.png screenshots")
    ap.add_argument("--out", required=True, help="output .pptx path")
    a = ap.parse_args()

    pngs = sorted(glob.glob(os.path.join(a.shots, "*.png")))
    if not pngs:
        raise SystemExit(f"no PNGs in {a.shots}")
    out = os.path.abspath(a.out)
    print(f"{len(pngs)} screenshots -> {out}")

    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    blank = prs.slide_layouts[BLANK_LAYOUT]

    for i, png in enumerate(pngs, start=1):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(os.path.abspath(png), Inches(0), Inches(0),
                                 width=Inches(W_IN), height=Inches(H_IN))
        print(f"  slide {i}: {os.path.basename(png)}")

    prs.save(out)
    print(f"SAVED -> {out}")

if __name__ == "__main__":
    main()
